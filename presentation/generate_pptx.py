"""Generate undata presentation as editable PPTX.

40-slide deck: light theme, teal/amber accents.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Colors — Light theme
BG = RGBColor(255, 255, 255)
TEXT = RGBColor(30, 41, 59)
ACCENT = RGBColor(14, 116, 144)       # Teal-700
ACCENT2 = RGBColor(180, 83, 9)        # Amber-700
MUTED = RGBColor(100, 116, 139)       # Slate-500
CARD = RGBColor(241, 245, 249)        # Slate-100
SUCCESS = RGBColor(21, 128, 61)       # Green-700
WARN = RGBColor(194, 65, 12)          # Orange-700
PURPLE = RGBColor(126, 34, 206)       # Purple-700

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]  # blank layout


def add_slide():
    slide = prs.slides.add_slide(BLANK)
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG
    # Bottom accent bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(7.45), prs.slide_width, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return slide


def add_text(slide, left, top, width, height, text, font_size=18, color=TEXT, bold=False, align=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = align
    return tf


def add_title(slide, text, y=0.4):
    tf = add_text(slide, 0.7, y, 11.5, 0.8, text, 36, TEXT, bold=True)
    # Underline bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(y + 0.7), Inches(2), Inches(0.04))
    bar.fill.solid()
    bar.fill.fore_color.rgb = ACCENT
    bar.line.fill.background()
    return y + 0.9


def add_bullets(slide, items, left=0.8, top=1.5, width=11, font_size=20, color=TEXT, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = spacing
        p.level = 0
        # Bullet
        pPr = p._pPr
        if pPr is None:
            from pptx.oxml.ns import qn
            pPr = p._p.get_or_add_pPr()
        from pptx.oxml.ns import qn
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        buClr = pPr.makeelement(qn('a:buClr'), {})
        srgb = buClr.makeelement(qn('a:srgbClr'), {'val': str(ACCENT)})
        buClr.append(srgb)
        # Remove existing bullets
        for old in pPr.findall(qn('a:buChar')):
            pPr.remove(old)
        for old in pPr.findall(qn('a:buClr')):
            pPr.remove(old)
        pPr.append(buClr)
        pPr.append(buChar)
    return tf


def add_card(slide, left, top, width, height, title, items, accent=ACCENT):
    # Card background
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = accent
    shape.line.width = Pt(2)
    # Title
    add_text(slide, left + 0.15, top + 0.1, width - 0.3, 0.4, title, 16, accent, bold=True)
    # Items
    y = top + 0.5
    for item in items:
        add_text(slide, left + 0.15, y, width - 0.3, 0.3, item, 13, MUTED)
        y += 0.25
    return shape


def add_stat(slide, left, top, number, label, color=ACCENT):
    add_text(slide, left, top, 2.5, 0.8, str(number), 48, color, bold=True)
    add_text(slide, left, top + 0.7, 2.5, 0.4, label, 16, MUTED)


# ============================================================
# SLIDES
# ============================================================

# 1 — Title
s = add_slide()
add_text(s, 0.7, 1.5, 10, 1.2, "undata", 72, ACCENT, bold=True)
add_text(s, 0.7, 2.7, 10, 0.7, "A Universal Data Element Registry", 40, TEXT, bold=True)
add_text(s, 0.7, 3.3, 10, 0.6, "for Neuroscience", 40, TEXT, bold=True)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.7), Inches(4.1), Inches(3), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT2; bar.line.fill.background()
add_text(s, 0.7, 4.3, 10, 0.5, "Building a shared language for neuroscience data schemas", 24, MUTED)
add_text(s, 0.7, 5.3, 5, 0.4, "Satrajit Ghosh", 22, TEXT)
add_text(s, 0.7, 5.7, 5, 0.4, "MIT  |  2026", 18, MUTED)

# 2 — Problem Hook
s = add_slide()
add_text(s, 2, 1.2, 10, 0.6, "Neuroscience has", 36, MUTED)
add_text(s, 2, 1.9, 10, 1, "no shared language", 64, ACCENT2, bold=True)
add_text(s, 2, 3.1, 10, 0.6, "for data elements.", 36, MUTED)
add_text(s, 0.7, 4.5, 12, 0.5, "5 major ecosystems  ×  5 different formats  ×  0 interoperability", 26, TEXT)
add_text(s, 0.7, 5.3, 12, 0.5, "BIDS  •  NWB  •  DANDI  •  openMINDS  •  AIND", 22, ACCENT)

# 3 — Three Problems
s = add_slide()
y = add_title(s, "Three Problems Nobody Has Solved")
add_card(s, 0.5, y, 3.8, 2.8, "1  SCHEMA FRAGMENTATION", [
    "Same concept, different formats",
    "JSON Schema, YAML, LinkML,",
    "Python classes, JSON-LD",
    "",
    'No way to ask: "What does BIDS',
    'call the thing NWB calls',
    'ElectrodeArray?"',
], ACCENT)
add_card(s, 4.6, y, 3.8, 2.8, "2  IDENTITY CONFLATION", [
    "What something IS ≠",
    "where it came FROM",
    "",
    "'age' in BIDS and 'age' in AIND",
    "are treated as different things",
    "despite being semantically",
    "identical",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 2.8, "3  INVISIBLE TRANSFORMS", [
    "Silent data transformations",
    "Unit conversions, type coercions",
    "happen in ad-hoc scripts",
    "",
    "No provenance record",
    "Breaking reproducibility",
], WARN)
add_text(s, 0.7, y + 3.1, 12, 0.4, "These problems scale with every new dataset, standard, and ecosystem.", 18, MUTED)

# 4 — Impact
s = add_slide()
y = add_title(s, "Why This Matters")
add_bullets(s, [
    "Researchers waste weeks mapping between formats manually",
    "Cross-ecosystem analyses require heroic bespoke engineering",
    "Transformations are undocumented → results not reproducible",
    "New standards fragment the landscape further",
    "No machine-readable crosswalk between any two ecosystems",
], top=y)
add_text(s, 0.7, 5.8, 12, 0.4, "The cost compounds with every dataset, every study, every collaboration.", 16, ACCENT2)

# 5 — Three Innovations
s = add_slide()
y = add_title(s, "Three Core Innovations")
add_card(s, 0.5, y, 3.8, 3.2, "CONTENT-ADDRESSED IDENTITY", [
    "SHA-256 hash from semantic",
    "properties (data type, unit,",
    "pattern, constraints, ontology)",
    "",
    "Two elements from different",
    "sources describing the same",
    "concept → same hash",
    "automatically",
], SUCCESS)
add_card(s, 4.6, y, 3.8, 3.2, "IDENTITY ≠ PROVENANCE", [
    "An element's identity (what it IS)",
    "is separate from provenance",
    "(where it came FROM)",
    "",
    "A single element accumulates",
    "provenance from multiple sources",
    "instead of duplicating",
], ACCENT)
add_card(s, 8.7, y, 3.8, 3.2, "EXPLICIT TRANSFORMS", [
    "Every conversion documented:",
    "• Unit (years → months)",
    "• Type (float → string)",
    "• Value (M/F → male/female)",
    "",
    "Content-addressed with their",
    "own provenance chains",
], ACCENT2)

# 6 — Pipeline
s = add_slide()
y = add_title(s, "The undata Pipeline")
stages = [("Extract", "Source adapters\n→ LinkML → entities", ACCENT),
          ("Enrich", "Embedding similarity\n+ LLM verification", SUCCESS),
          ("Align", "Cross-source\nalias detection", ACCENT2),
          ("Commit", "Content addressing\n+ dedup + merge", WARN),
          ("Transform", "Conversion logic\ngeneration", PURPLE)]
for i, (name, desc, color) in enumerate(stages):
    left = 0.5 + i * 2.5
    add_card(s, left, y, 2.2, 1.8, name, desc.split("\n"), color)
    if i < 4:
        add_text(s, left + 2.25, y + 0.7, 0.3, 0.4, "→", 24, MUTED)

add_text(s, 0.7, y + 2.2, 12, 0.3, "Sources:", 16, MUTED, bold=True)
add_text(s, 0.7, y + 2.5, 12, 0.3, "BIDS • NWB • DANDI • openMINDS • AIND • OpenNeuro • ReproSchema • NDA", 15, TEXT)
add_text(s, 0.7, y + 3.0, 12, 0.3, "Knowledge Service:", 16, MUTED, bold=True)
add_text(s, 0.7, y + 3.3, 12, 0.3, "13 ontologies  •  2.99M terms  •  268K embedded  •  pyoxigraph RDF store", 15, ACCENT)

# 7 — Architecture
s = add_slide()
y = add_title(s, "Three-Layer Architecture")
add_card(s, 0.5, y, 3.8, 2.5, "FRONTEND — Next.js", [
    "Apollo Client + Tailwind CSS",
    "Browse, search, curation, chat",
    "Entity navigator (ER diagram)",
    "Infinite scroll + server-side sort",
    "Real-time LLM curation assistant",
], ACCENT)
add_card(s, 4.6, y, 3.8, 2.5, "BACKEND — FastAPI", [
    "Strawberry GraphQL + PostgreSQL 16",
    "pgvector for embeddings",
    "JSONB for flexible entity storage",
    "Keycloak OIDC authentication",
    "PROV-O audit trail for all mutations",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 2.5, "LIBRARY — Pure Python", [
    "No DB dependencies, standalone CLI",
    "All pipeline logic: extract → transform",
    "StorageBackend protocol (file or DB)",
    "436 tests, Python 3.14",
    "uv for dependency management",
], SUCCESS)
add_text(s, 0.7, y + 2.8, 12, 0.3, "StorageBackend Protocol:  read() • write() • list() • exists() • delete() • merge_provenance()", 16, TEXT)
add_text(s, 0.7, y + 3.2, 12, 0.3, "FileBackend (YAML + Parquet + pyoxigraph)  |  DatabaseBackend (PostgreSQL + pgvector)", 15, MUTED)

# 8 — Entity Model
s = add_slide()
y = add_title(s, "Content-Addressed Entity Model")
types = [("Element", "Field-level: age, weight\ntype, unit, constraints", ACCENT),
         ("Schema", "Class-level: Subject, Session\nproperties, inheritance", ACCENT2),
         ("Value", "Categorical: male, MRI\nontology grounding", SUCCESS),
         ("ValueSet", "Collection: modalities\nmembers list", WARN)]
for i, (name, desc, color) in enumerate(types):
    left = 0.5 + i * 3.1
    add_card(s, left, y, 2.8, 2.8, name, [
        *desc.split("\n"), "",
        "sha256 (identity)",
        "semantic{} (content)",
        "provenance[] (origins)",
        "annotations[] (ontology)",
    ], color)
add_text(s, 0.7, y + 3.1, 12, 0.3, "+ Transform (conversions)  + CurationFlag (quality)  + AuditLog (provenance)", 18, ACCENT)
add_text(s, 0.7, y + 3.5, 12, 0.3, "All entities: staged → curated through curator review. No automated changes without evidence.", 16, MUTED)

# 9 — Enrichment
s = add_slide()
y = add_title(s, "Multi-Precision Enrichment")
add_bullets(s, [
    "Embedding similarity: all-MiniLM-L6-v2 → cosine similarity to 268K ontology terms",
    "SKOS mapping: exactMatch (≥0.95) • closeMatch (≥0.85) • broadMatch • relatedMatch",
    "LLM verification for borderline matches (0.5–0.7) via litellm",
    "Ontology hierarchy traversal: ancestors added as broadMatch",
    "Evidence Chain for every annotation: score + URI verification + reasoning",
    "Curated annotations protected from re-enrichment",
], top=y, font_size=20)
add_text(s, 0.7, 5.8, 12, 0.3, "NCBITaxon filtered to 20 neuroscience-relevant species (from 2.7M terms)", 16, ACCENT2)

# 10 — Speckit Process
s = add_slide()
y = add_title(s, "Engineering Process: Speckit Lifecycle")
stages = [("/speckit.specify", "User stories\nAcceptance criteria"),
          ("/speckit.clarify", "Ambiguity detection\nSpec updates"),
          ("/speckit.plan", "Architecture\nData model, contracts"),
          ("/speckit.tasks", "Dependency-ordered\ntask breakdown"),
          ("/speckit.analyze", "Cross-artifact\nconsistency check"),
          ("/speckit.implement", "Phase-by-phase\nexecution")]
for i, (name, desc) in enumerate(stages):
    left = 0.3 + i * 2.1
    color = ACCENT if i % 2 == 0 else ACCENT2
    add_card(s, left, y, 1.9, 1.8, name, desc.split("\n"), color)
    if i < 5:
        add_text(s, left + 1.95, y + 0.7, 0.2, 0.3, "→", 18, MUTED)
add_text(s, 0.7, y + 2.2, 12, 0.3, "38 features specified and implemented through this lifecycle", 20, TEXT)
add_text(s, 0.7, y + 2.6, 12, 0.3, "Project constitution (v2.1.0) with 7 core principles enforced at every stage", 16, MUTED)

# 11 — Constitution
s = add_slide()
y = add_title(s, "Project Constitution: 7 Core Principles")
principles = [
    ("I. Simplicity First", "YAGNI at all levels — complexity must be justified"),
    ("II. Test-Driven Development", "Red-Green-Refactor, non-negotiable"),
    ("III. API-First Design", "Contracts before implementation"),
    ("IV. Observability", "Structured JSON logs, no silent failures"),
    ("V. No Deprecation", "System has no consumers — rewrite freely"),
    ("VI. Environment Isolation", "uv + venvs, bridge venv exception for adapters"),
    ("VII. Developer Experience", "docker compose up → everything works"),
]
for i, (title, desc) in enumerate(principles):
    color = ACCENT if i % 2 == 0 else ACCENT2
    cy = y + i * 0.65
    add_text(s, 0.8, cy, 4, 0.35, title, 18, color, bold=True)
    add_text(s, 5.2, cy, 8, 0.35, desc, 18, MUTED)
add_text(s, 0.7, y + 4.8, 12, 0.3, "Version 2.1.0  •  Ratified 2026-03-07  •  Enforced on every PR", 14, MUTED)

# 12 — Feature Evolution
s = add_slide()
y = add_title(s, "38 Features: From Brainstorm to Production")
add_text(s, 0.8, y, 12, 0.35, "Iteration 1: Exploration (001–027)", 20, MUTED, bold=True)
add_text(s, 0.8, y+0.4, 12, 0.3, "001–010: Schema integration, backend, explorer, migration, enrichment", 16, MUTED)
add_text(s, 0.8, y+0.7, 12, 0.3, "011–020: Metamodel, provenance, fullstack, ingestion overhaul, extraction", 16, MUTED)
add_text(s, 0.8, y+1.0, 12, 0.3, "021–027: Source acquisition, repo restructure, ontology, staged enrichment", 16, MUTED)
add_text(s, 0.8, y+1.5, 12, 0.35, "Iteration 2: Production (028–038)", 20, ACCENT, bold=True)
add_text(s, 0.8, y+1.9, 12, 0.3, "028–029: Storage abstraction, backend service layer", 16, TEXT)
add_text(s, 0.8, y+2.2, 12, 0.3, "030–032: Frontend integration, CivicDB-inspired UI, authentication", 16, TEXT)
add_text(s, 0.8, y+2.5, 12, 0.3, "033–034: Unit standardization, curation interface", 16, TEXT)
add_text(s, 0.8, y+2.8, 12, 0.3, "035–037: UX overhaul, knowledge service, data export", 16, TEXT)
add_text(s, 0.8, y+3.1, 12, 0.3, "038: System hardening — 47 tasks, 9 user stories, all complete", 16, ACCENT2)

# 13 — Registry Scale
s = add_slide()
y = add_title(s, "The Registry: Scale & Coverage")
add_stat(s, 0.7, y, "7,055", "Elements", ACCENT)
add_stat(s, 3.2, y, "1,005", "Schemas", ACCENT2)
add_stat(s, 5.7, y, "15,881", "Values", SUCCESS)
add_stat(s, 8.2, y, "2,426", "Value Sets", WARN)
add_stat(s, 10.7, y, "15", "Transforms", PURPLE)
sources = [("BIDS", "585 elements • JSON Schema"), ("NWB", "179 elements • Python/hdmf"),
           ("DANDI", "398 elements • Pydantic"), ("openMINDS", "473 elements • JSON-LD"),
           ("AIND", "556 elements • JSON Schema"), ("OpenNeuro", "28 elements • TSV/datalad"),
           ("ReproSchema", "4,836 elements • JSON-LD")]
add_text(s, 0.7, y+1.8, 12, 0.3, "7 Sources:", 18, TEXT, bold=True)
for i, (name, detail) in enumerate(sources):
    col, row = i % 4, i // 4
    add_text(s, 0.8 + col*3.1, y+2.2+row*0.6, 3, 0.3, name, 16, ACCENT, bold=True)
    add_text(s, 0.8 + col*3.1, y+2.5+row*0.6, 3, 0.3, detail, 14, MUTED)

# 14 — Enrichment Coverage
s = add_slide()
y = add_title(s, "Enrichment Coverage")
sources_e = [("openMINDS", "70%", SUCCESS), ("AIND", "34.4%", ACCENT), ("DANDI", "25.9%", ACCENT),
             ("BIDS", "10.3%", ACCENT2), ("NWB", "1.1%", WARN)]
for i, (name, pct, color) in enumerate(sources_e):
    cy = y + i * 0.55
    add_text(s, 0.8, cy, 2, 0.35, name, 20, TEXT)
    add_text(s, 3.5, cy, 2, 0.35, pct, 20, color, bold=True)
add_text(s, 0.8, y + 3.0, 12, 0.3, "13 Ontologies  •  2.99M total terms  •  268K embedded", 20, ACCENT)
add_bullets(s, [
    "Evidence Chain: similarity score + URI verification + reasoning chain",
    "LLM verification (gpt-5.4-nano: 0.06s/pair, ollama/qwen3.5: 0.5s/pair)",
    "NCBITaxon filtered to 20 neuroscience species (from 2.7M terms)",
], top=y + 3.5, font_size=18)

# 15 — Knowledge Service
s = add_slide()
y = add_title(s, "Knowledge Service: Expanding the Semantic Universe")
add_card(s, 0.5, y, 3.8, 3, "ONTOLOGIES", [
    "Core: NCIT, PATO, UBERON, HP, OBI",
    "Neuro: NIDM, CogPO, NIF-std",
    "Imaging: DICOM, RadLex, NeuroNames",
    "New: HoMBA (brain anatomy)",
    "",
    "pyoxigraph RDF + SPARQL",
    "Checksum-based dedup per source",
], ACCENT)
add_card(s, 4.6, y, 3.8, 3, "DATA REPOSITORIES", [
    "OpenNeuro (BIDS, datalad API)",
    "DANDI (NWB + schema.org)",
    "openMINDS instances (4,390)",
    "NDA data dictionary API",
    "ReproSchema library (JSON-LD)",
    "",
    "Adapter pattern: BaseAdapter",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 3, "KNOWLEDGE BASES", [
    "Cognitive Atlas (concepts→tasks)",
    "InterLex (NIF registry)",
    "NITRC (tools + capabilities)",
    "NIH CDE Repository",
    "CDISC, FHIR, schema.org",
    "",
    "Discovery service polls daily",
], SUCCESS)

# 16 — Adapters
s = add_slide()
y = add_title(s, "Adapter Architecture: 8 Source Adapters")
adapters = [("BIDS", "LinkML-first, raw JSON Schema descriptors"),
            ("NWB", "Python/hdmf introspection via bridge venv (Python 3.12)"),
            ("DANDI", "Pydantic model introspection, inheritance handling"),
            ("openMINDS", "JSON-LD → LinkML, 4,390 instance values"),
            ("AIND", "JSON Schema with $ref resolution, enum extraction"),
            ("OpenNeuro", "datalad Python API, TSV/CSV + JSON sidecar scanning"),
            ("ReproSchema", "JSON-LD with relative path ref resolution"),
            ("NDA", "REST API, data dictionary extraction with value ranges")]
for i, (name, desc) in enumerate(adapters):
    cy = y + i * 0.5
    add_text(s, 0.8, cy, 2.5, 0.35, name, 18, ACCENT, bold=True)
    add_text(s, 3.5, cy, 9, 0.35, desc, 17, TEXT)
add_text(s, 0.7, y + 4.3, 12, 0.3, "All adapters: BaseAdapter → extract() → list[ClassifiedEntity]", 16, MUTED)

# 17 — Content Addressing
s = add_slide()
y = add_title(s, "Content-Addressed Identity: How It Works")
add_card(s, 0.5, y, 5.8, 3.5, "TWO-MODE HASHING", [
    "Mode 1 — Ontology-Anchored:",
    "  hash(data_type + unit + pattern + constraints",
    "       + primary_ontology_uri)",
    "  → Same concept across sources gets same hash",
    "",
    "Mode 2 — Structural Fallback:",
    "  hash(data_type + unit + pattern + constraints",
    "       + class + attribute + description)",
    "  → Unique within context when no ontology match",
    "",
    "Short key: 12 hex chars from SHA-256",
], ACCENT)
add_card(s, 6.6, y, 5.8, 3.5, "PROVENANCE ACCUMULATION", [
    "Element: age (sha256: 240539a4c71f...)",
    "",
    "Provenance[0]: source=bids, class=participants",
    "Provenance[1]: source=nwb, class=Subject",
    "Provenance[2]: source=aind, class=Subject",
    "Provenance[3]: source=openneuro/ds000228",
    "",
    "One entity, many origins",
    "Cross-referenced automatically",
], ACCENT2)

# 18 — Transforms
s = add_slide()
y = add_title(s, "Transform Generation: Three Strategies")
add_card(s, 0.5, y, 3.8, 3, "1. SHARED ONTOLOGY URI", [
    "Elements with same primary",
    "annotation (e.g., NCIT:C25150)",
    "Original method: 15 transforms",
    "",
    "Detects: identity, unit_conversion",
    "type_conversion, value_mapping",
    "structural, unknown",
], ACCENT)
add_card(s, 4.6, y, 3.8, 3, "2. NAME-BASED MATCHING", [
    "Group by provenance name",
    "(case-insensitive)",
    "Cross-source only",
    "(BIDS age ↔ NWB age)",
    "",
    "Dramatically increases coverage:",
    "100+ transforms (up from 15)",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 3, "3. EMBEDDING SIMILARITY", [
    "Cosine similarity between",
    "element description embeddings",
    "Threshold: 0.8 (cross-source)",
    "",
    "Catches semantic matches missed",
    "by name or ontology matching",
    "Capped at 500 pairs",
], SUCCESS)
add_text(s, 0.7, y+3.2, 12, 0.3, "Many-to-one: source_elements[] for composite transforms (age + age_unit → age_in_years)", 16, MUTED)

# 19 — Curation
s = add_slide()
y = add_title(s, "Evidence-Based Curation Workflow")
add_card(s, 0.5, y, 3.8, 3.5, "EVIDENCE CHAIN", [
    "similarity_score: 0.85",
    "similarity_method: cosine",
    "source_text: 'Age in years'",
    "target: NCIT:C25150 'Age'",
    "uri_verified: true ✓",
    "reasoning: 'Cosine similarity...'",
    "",
    "No hallucinated confidence.",
    "Every claim backed by",
    "verifiable evidence.",
], ACCENT)
add_card(s, 4.6, y, 3.8, 3.5, "CURATION ROLES", [
    "Contributor (default):",
    "  Suggest annotations, flag issues",
    "",
    "Curator (reviews):",
    "  Approve/reject proposals",
    "  Resolve curation flags",
    "  LLM-assisted chat curation",
    "",
    "Admin (manages):",
    "  Import, manage ontologies",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 3.5, "AUDIT TRAIL", [
    "PROV-O style for every mutation:",
    "  agent (who)",
    "  activity (what)",
    "  entity (on what)",
    "  generated_entity (new version)",
    "  timestamp",
    "",
    "Queryable by entity, user, time",
    "14,984 curation flags tracked",
], SUCCESS)

# 20 — LLM Chat
s = add_slide()
y = add_title(s, "LLM-Powered Curation Assistant")
add_bullets(s, [
    "Real-time chat connected to configurable LLM (litellm: OpenAI, ollama, Anthropic)",
    "Entity context injected into every conversation (fields, provenance, annotations)",
    "Auto-suggest improvements on entity load (missing annotations, units, descriptions)",
    "Tool execution: lookup_ontology_term, propose_entity_change, fetch_entity",
    "Proposals appear as reviewable diffs with evidence chains",
    "Borderline candidates (0.5–0.7) batch-verified via LLM with disk cache",
], top=y, font_size=20)

# 21 — Search
s = add_slide()
y = add_title(s, "Search: Lexical, Semantic, and Combined")
add_card(s, 0.5, y, 3.8, 2.5, "LEXICAL SEARCH", [
    "PostgreSQL tsvector full-text",
    "ILIKE fallback for simple queries",
    "Cross-entity: elements, schemas,",
    "values, valuesets",
], ACCENT)
add_card(s, 4.6, y, 3.8, 2.5, "SEMANTIC SEARCH", [
    "pgvector nearest-neighbor",
    "Query → embedding → cosine dist",
    '"brain area" → brain_region,',
    "cortical_area, anatomical_region",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 2.5, "COMBINED MODE", [
    "Default: lexical + semantic",
    "Lexical matches first (exact)",
    "Semantic matches after (related)",
    "Mode toggle in search page UI",
], SUCCESS)

# 22 — UI Browse
s = add_slide()
y = add_title(s, "UI: Entity Navigator — ER Diagram in a Browser")
add_bullets(s, [
    "Every entity links to related entities — schemas ↔ elements ↔ values ↔ transforms",
    "Server-side sorting on all columns (not client-side re-sort of loaded page)",
    "Infinite scroll with IntersectionObserver sentinel (200px before end)",
    "Source filtering on every browse page (BIDS, NWB, DANDI, openMINDS, AIND, ...)",
    "Cross-entity search with results grouped by type",
    "Annotation chips with expandable evidence chains",
    "Provenance badge strips showing all source attestations",
], top=y, font_size=20)

# 23 — UI Detail
s = add_slide()
y = add_title(s, "UI: Entity Detail — Full Semantic Context")
add_bullets(s, [
    "Identity: sha256, data type, unit, unit URI, pattern, min/max, value domain",
    "Semantic fields: description, question text, response options",
    "Ontology annotations with SKOS relation, score, expandable evidence chain",
    "Provenance entries with W3C PROV-O metadata (source, class, activity, timestamp)",
    "Cross-references: schemas containing this element, transforms involving it",
    "Curation flags with context, LLM verification, resolve buttons",
    "Versioning: superseded_by link chain, curation_update transforms",
], top=y, font_size=20)

# 24 — Feature 038
s = add_slide()
y = add_title(s, "Feature 038: System Hardening")
add_text(s, 0.8, y, 12, 0.4, "47 tasks  •  9 user stories  •  12 phases  •  47/47 complete", 22, ACCENT2)
stories = [("US1 P1", "LLM Curation Chat", "Auto-suggest, evidence chains, tool execution"),
           ("US2 P1", "Name-Based Transforms", "100+ transforms (up from 15), many-to-one"),
           ("US3 P1", "Additional Sources", "OpenNeuro, ReproSchema, NDA adapters"),
           ("US4 P2", "Search Modes", "Lexical/semantic/both with pgvector"),
           ("US5 P2", "Ontology Admin", "pyoxigraph store info, NCBITaxon filter"),
           ("US6 P2", "Server-Side Sorting", "All browse pages, infinite scroll"),
           ("US7 P3", "Audit + Downloads", "PROV-O audit, nightly exports"),
           ("US8 P3", "CI + Pipeline", "Action v6, index auto-rebuild, LLM enrichment"),
           ("US9 P1", "Versioned Dependencies", "Checksum detection, auto re-enrich")]
for i, (us, title, desc) in enumerate(stories):
    color = ACCENT2 if "P1" in us else ACCENT if "P2" in us else MUTED
    cy = y + 0.5 + i * 0.45
    add_text(s, 0.8, cy, 1.5, 0.35, us, 16, color, bold=True)
    add_text(s, 2.5, cy, 3.5, 0.35, title, 16, TEXT)
    add_text(s, 6.2, cy, 6.5, 0.35, desc, 15, MUTED)

# 25 — Evidence Chain Detail
s = add_slide()
y = add_title(s, "Evidence Chain: No Hallucinated Confidence")
add_card(s, 0.5, y, 5.8, 3.5, "EVIDENCE CHAIN STRUCTURE", [
    "┌─────────────────────────────────────┐",
    "│  Similarity Score: 0.85  [85%]         │",
    "│  Method: cosine_embedding               │",
    "│  Source: 'Age of participant (years)'    │",
    "│  Target: NCIT:C25150 'Age'              │",
    "│  Def: 'How long something existed'      │",
    "│  URI Verified: ✓ (HTTP 200)             │",
    "│  Reasoning: Cosine similarity 0.85      │",
    "│  between element and term. Relation:    │",
    "│  skos:closeMatch                        │",
    "└─────────────────────────────────────┘",
], ACCENT)
add_card(s, 6.6, y, 5.8, 3.5, "THREE VERIFICATION LAYERS", [
    "1. SEMANTIC SIMILARITY",
    "   Embedding cosine score (0.0–1.0)",
    "   Threshold: 0.7 for annotation",
    "",
    "2. LINK VERIFICATION",
    "   HTTP HEAD check on proposed URI",
    "   Detects stale ontology references",
    "",
    "3. REASONING CHAIN",
    "   Step-by-step explanation",
    "   LLM or template generated",
], ACCENT2)

# 26 — Version Management
s = add_slide()
y = add_title(s, "Versioned Dependency Management")
add_bullets(s, [
    "Scheduled checksum comparison for all registered ontologies and sources",
    "When change detected → automatic re-enrichment of affected entities",
    "Curator-approved annotations preserved (curated_annotations field)",
    "Only automated annotations from old version are re-evaluated",
    "VersionTransition recorded: old checksum → new checksum → timestamp",
    "checkDependencyVersions GraphQL mutation for manual trigger",
    "HoMBA ontology loaded from brain-bican GitHub releases (OWL→TTL fallback)",
], top=y, font_size=20)

# 27 — Polyglot Storage
s = add_slide()
y = add_title(s, "Polyglot Storage: Right Tool for Each Data Type")
rows = [("Entities", "CRUD, filter, paginate", "YAML", "PostgreSQL JSONB"),
        ("Ontology terms", "Graph traversal, SPARQL", "pyoxigraph", "RDF store"),
        ("Embeddings", "Nearest-neighbor, cosine", "Parquet + numpy", "pgvector"),
        ("Full-text search", "Keyword, faceted", "grep/memory", "tsvector"),
        ("LLM cache", "Key-value lookup", "JSON file", "PostgreSQL"),
        ("Curation flags", "CRUD, status filter", "YAML", "PostgreSQL"),
        ("Audit log", "Append-only, time-range", "YAML", "PostgreSQL")]
headers = ["Data", "Pattern", "File Backend", "DB Backend"]
for i, h in enumerate(headers):
    add_text(s, 0.8+i*3, y, 2.8, 0.35, h, 16, ACCENT, bold=True)
for j, row in enumerate(rows):
    for i, cell in enumerate(row):
        add_text(s, 0.8+i*3, y+0.5+j*0.42, 2.8, 0.35, cell, 15, TEXT if i == 0 else MUTED)

# 28 — GraphQL API
s = add_slide()
y = add_title(s, "GraphQL API: Single Unified Interface")
add_card(s, 0.5, y, 5.8, 3.5, "QUERIES", [
    "element(sha256) → Element",
    "browseElements(source, sort, search)",
    "browseSchemas / browseValues / browseValuesets",
    "search(query, mode) → SearchResult[]",
    "curationQueue(flagType, status)",
    "flagsForEntity / schemasUsingElement",
    "transformsForElement / ontologyStoreInfo",
    "auditLog(entity, agent, activity)",
    "enrichmentProposals / releases",
], ACCENT)
add_card(s, 6.6, y, 5.8, 3.5, "MUTATIONS", [
    "resolveFlag / batchResolveFlags",
    "updateElement / updateSchema / updateValue",
    "approveAnnotation / rejectAnnotation",
    "versionElement(sha256, changes)",
    "approveIngestion / rejectIngestion",
    "reviewProposal / requestEnrichment",
    "importRegistry / exportRegistry",
    "checkDependencyVersions / tagRelease",
], ACCENT2)

# 29 — Testing
s = add_slide()
y = add_title(s, "Testing: 436 Library Tests + Full Stack CI")
add_stat(s, 0.7, y, "436", "Library Tests", SUCCESS)
add_stat(s, 3.2, y, "16", "Transform Tests", ACCENT)
add_stat(s, 5.7, y, "9", "CI Workflows", ACCENT2)
add_stat(s, 8.2, y, "0", "TS Errors", SUCCESS)
add_bullets(s, [
    "pytest + pytest-asyncio for all Python tests",
    "TypeScript strict mode with zero non-test errors",
    "GitHub Actions: library, backend, frontend, e2e, lint, build-images",
    "Actions updated to v5/v6 (Node.js 24 compatible)",
    "Docker compose for integration testing",
], top=y + 1.8, font_size=18)

# 30 — Prior Work
s = add_slide()
y = add_title(s, "Standing on the Shoulders of Giants")
add_card(s, 0.5, y, 3.8, 3.5, "NEUROSCIENCE STANDARDS", [
    "BIDS — Gorgolewski et al., 2016",
    "NWB — Teeters et al., 2015",
    "DANDI — Halchenko et al., 2024",
    "openMINDS — Zehl et al., 2023",
    "AIND — Allen Institute, 2023",
    "ReproSchema — ReproNim",
    "NDA — NIMH Data Archive",
], ACCENT)
add_card(s, 4.6, y, 3.8, 3.5, "KNOWLEDGE RESOURCES", [
    "NCIT — NCI Thesaurus (209K)",
    "UBERON — Anatomy ontology",
    "PATO — Phenotype/trait",
    "RadLex — Radiology (46K)",
    "DICOM — Imaging (5K)",
    "NIDM — Neuroimaging DM",
    "HoMBA — Brain anatomy",
    "Cognitive Atlas — Poldrack et al.",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 3.5, "TECHNOLOGIES", [
    "LinkML — Moxon et al., 2021",
    "SKOS — W3C vocabulary",
    "W3C PROV-O — provenance",
    "sentence-transformers",
    "pyoxigraph — RDF store",
    "pgvector — embeddings",
    "litellm — LLM abstraction",
    "datalad — distributed data",
    "CivicDB — curation model",
], SUCCESS)

# 31 — Community Model
s = add_slide()
y = add_title(s, "Community Curation Model (Inspired by CivicDB)")
add_bullets(s, [
    "Three-tier graduated engagement: Contributor → Curator → Admin",
    "Revision-based: suggestions with evidence, curator approves/rejects",
    "Polymorphic concerns: any entity type can be commented, flagged, subscribed to",
    "Evidence panels: automated match candidates + scores, LLM verification",
    "Activity feed filterable by action, user, source",
    "Organization-level attribution and contribution statistics",
], top=y, font_size=20)
add_text(s, 0.7, 5.5, 12, 0.3, "Every change goes through curation — no automated change applied without review.", 18, ACCENT2)

# 32 — Exports
s = add_slide()
y = add_title(s, "Data Export & Downloads")
add_bullets(s, [
    "Nightly export scheduler: asyncio background task, daily archive production",
    "Export format: YAML entities + Parquet embeddings + manifest",
    "Release records: version, file size, entity counts, download count",
    "Static file serving at /api/downloads/",
    "Downloads page in UI with version, date, size, entity counts",
    "Versioned releases: tag any nightly export with a version string",
], top=y, font_size=20)

# 33 — Cross-Ecosystem
s = add_slide()
y = add_title(s, "Cross-Ecosystem Impact")
add_text(s, 0.8, y, 5, 0.35, "Before undata:", 20, WARN, bold=True)
add_bullets(s, [
    "Each ecosystem defines 'age' independently",
    "Converting BIDS↔NWB requires custom scripts per field",
    "No machine-readable way to discover equivalences",
    "Transformation provenance lost in pipeline scripts",
], top=y+0.4, font_size=18)
add_text(s, 0.8, y+2.5, 5, 0.35, "After undata:", 20, SUCCESS, bold=True)
add_bullets(s, [
    "One 'age' element with provenance from 5+ sources",
    "Transforms auto-generated with documented logic",
    "Semantic search finds related concepts across ecosystems",
    "Every annotation backed by evidence chain + curation",
], top=y+2.9, font_size=18)

# 34 — Reproducibility
s = add_slide()
y = add_title(s, "Reproducibility Through Provenance")
add_bullets(s, [
    "Content-addressed identity: same semantic → same hash, deterministic",
    "W3C PROV-O audit trail: who, what, when, why, and what was generated",
    "Version transitions: old ontology → new ontology → affected entities",
    "Transform provenance: explicit conversion logic, not hidden scripts",
    "Curation flags: machine-generated quality concerns for human review",
    "Nightly exports: timestamped, versioned snapshots of complete registry",
], top=y, font_size=20)

# 35 — Scalability
s = add_slide()
y = add_title(s, "Scalability & Performance")
add_stat(s, 0.7, y, "<5s", "Chat response", ACCENT)
add_stat(s, 3.2, y, "<1s", "Search latency", ACCENT2)
add_stat(s, 5.7, y, "<10m", "Pipeline run", SUCCESS)
add_stat(s, 8.2, y, "7K+", "Elements", WARN)
add_bullets(s, [
    "PostgreSQL 16 JSONB for flexible entity storage at scale",
    "pgvector for embedding nearest-neighbor in milliseconds",
    "Cursor-based pagination (not offset) for consistent performance",
    "Ontology store handles 2.99M terms in pyoxigraph",
    "Apollo cache with paginationMerge for frontend efficiency",
], top=y+1.8, font_size=18)

# 36 — Developer Experience
s = add_slide()
y = add_title(s, "Developer Experience: One Command to Everything")
add_text(s, 0.8, y, 12, 0.5, "$ docker compose up -d", 28, ACCENT)
add_text(s, 0.8, y+0.5, 12, 0.4, "→ Database seeded  •  Backend running  •  Frontend live  •  Auth configured", 18, MUTED)
add_bullets(s, [
    "Library standalone: uv run undata-library pipeline --source bids (no Docker)",
    "Hot reload for backend and frontend development",
    "docker-compose.override.yml for full registry mounting",
    "Curated seed subset (70 elements, 7 sources) committed for quick start",
    "Full registry at ~/.cache/undata/registry/ for comprehensive testing",
], top=y+1.2, font_size=18)

# 37 — What's Next
s = add_slide()
y = add_title(s, "What's Next")
add_card(s, 0.5, y, 3.8, 3, "NEAR TERM", [
    "Full enrichment with all 7 sources",
    "100+ transforms (3 strategies)",
    "Deploy to staging environment",
    "Community beta with curators",
    "HoMBA + additional ontologies",
    "OpenNeuro batch (100+ datasets)",
], ACCENT)
add_card(s, 4.6, y, 3.8, 3, "MEDIUM TERM", [
    "Custom LLM fine-tuning",
    "Real-time collaborative curation",
    "Cross-registry federation",
    "API consumers and integrations",
    "Automated pipeline scheduling",
], ACCENT2)
add_card(s, 8.7, y, 3.8, 3, "LONG TERM", [
    "Community-driven ontology creation",
    "Federated registry network",
    "Standard adoption by ecosystems",
    "Training data for neuro NLP",
    "Integration with analysis pipelines",
], SUCCESS)

# 38 — Key Takeaways
s = add_slide()
add_text(s, 2, 0.5, 10, 0.8, "Key Takeaways", 54, ACCENT, bold=True)
add_bullets(s, [
    "Content-addressed identity solves schema fragmentation across 7+ ecosystems",
    "Evidence-based curation prevents hallucinated confidence — every claim verifiable",
    "38 features, 436 tests, 7,055 elements from a rigorous speckit-driven process",
    "Three enrichment strategies: embedding + LLM verification + ontology hierarchy",
    "Three transform strategies: shared URI + name matching + embedding similarity",
    "Full audit trail with PROV-O provenance for every mutation",
    "Library + backend + frontend — all open, all tested, all documented",
], top=1.5, font_size=20)

# 39 — By The Numbers
s = add_slide()
y = add_title(s, "By The Numbers")
stats = [("38", "Features specified", ACCENT), ("47", "Tasks in hardening", ACCENT2),
         ("436", "Library tests", SUCCESS), ("7,055", "Elements", ACCENT),
         ("15,881", "Values", ACCENT2), ("1,005", "Schemas", SUCCESS),
         ("2,426", "Value Sets", WARN), ("9", "Ontologies loaded", ACCENT),
         ("268K", "Terms embedded", ACCENT2), ("8", "Source adapters", SUCCESS),
         ("7", "Principles", WARN), ("3", "Architecture layers", ACCENT)]
for i, (num, label, color) in enumerate(stats):
    col, row = i % 4, i // 4
    add_stat(s, 0.7 + col * 3.1, y + row * 1.5, num, label, color)

# 40 — Thank You
s = add_slide()
add_text(s, 2, 1.5, 10, 1, "Thank You", 72, ACCENT, bold=True, align=PP_ALIGN.CENTER)
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(2.7), Inches(2.5), Inches(0.04))
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT2; bar.line.fill.background()
add_text(s, 2, 3, 10, 0.5, "undata: A Universal Data Element Registry for Neuroscience", 24, TEXT, align=PP_ALIGN.CENTER)
add_text(s, 2, 4, 10, 0.5, "github.com/sensein/undata", 22, ACCENT, align=PP_ALIGN.CENTER)
add_text(s, 2, 4.5, 10, 0.4, "satra@mit.edu", 20, MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1, 5.5, 11, 0.4, "Built with: Python 3.14 • FastAPI • Next.js • PostgreSQL • pyoxigraph", 16, MUTED, align=PP_ALIGN.CENTER)
add_text(s, 1, 5.9, 11, 0.4, "Powered by: speckit lifecycle • litellm • sentence-transformers • datalad", 16, MUTED, align=PP_ALIGN.CENTER)

# ============================================================
# APPENDIX — Feature Groupings + Architecture Diagrams
# ============================================================

# A1 — Feature Groupings: Iteration 1
s = add_slide()
y = add_title(s, "Appendix: Feature Map — Iteration 1 (001–027)")
groups = [
    ("Foundation (001–005)", ["001 Schema integration", "002 Schema backend", "003 Schema explorer", "004 Migration API", "005 Schema enrichment"]),
    ("Adapters (006–010)", ["006 Dual-path adapters", "007 End-to-end pipeline", "008 Schema import roundtrip", "009 Tutorials", "010 JupyterBook"]),
    ("Core Model (011–015)", ["011 Metamodel provenance", "012 Fullstack compose", "013 Migration UI", "014 Deployment pipeline", "015 undata-library"]),
    ("Data Richness (016–020)", ["016 Value concepts", "017 Backend-library alignment", "018 Rich data model", "019 Ingestion overhaul", "020 Extract transforms"]),
    ("Pipeline (021–027)", ["021 Source acquisition", "022 Full re-extract", "023 Repo restructure", "024 Ontology service", "025 Ontology expansion", "026 Staged enrichment", "027 Library hardening"]),
]
for i, (group, features) in enumerate(groups):
    col, row = i % 3, i // 3
    left = 0.5 + col * 4.2
    top = y + row * 2.8
    add_card(s, left, top, 3.9, 2.5, group, features, MUTED if row == 0 and col < 2 else ACCENT if row == 0 else ACCENT2)

# A2 — Feature Groupings: Iteration 2
s = add_slide()
y = add_title(s, "Appendix: Feature Map — Iteration 2 (028–038)")
groups2 = [
    ("Infrastructure (028–029)", ["028 Storage abstraction — StorageBackend protocol", "029 Backend service — FastAPI + Strawberry GraphQL"]),
    ("Frontend (030–032)", ["030 Frontend integration — Next.js + Apollo", "031 CivicDB UI — community curation model", "032 Authentication — Keycloak OIDC"]),
    ("Data Quality (033–034)", ["033 Unit standardization — QUDT + cmixf", "034 Curation interface — flags, queue, chat"]),
    ("UX + Knowledge (035–037)", ["035 UX overhaul — EntityDataGrid, PropertyTable, search", "036 Knowledge service — ontology store, adapters, LLM", "037 Data export — YAML + Parquet + manifest"]),
    ("Hardening (038)", ["038 System hardening — 47 tasks, 9 user stories:", "  LLM chat, evidence chains, transforms,", "  search modes, audit log, versioning,", "  7 sources, downloads, CI updates"]),
]
for i, (group, features) in enumerate(groups2):
    col, row = i % 3, i // 3
    left = 0.5 + col * 4.2
    top = y + row * 2.8
    add_card(s, left, top, 3.9, 2.5, group, features, ACCENT if i < 2 else ACCENT2 if i < 4 else RGBColor(21, 128, 61))

# A3-A7 — Architecture diagrams as image slides
from pathlib import Path
diagram_dir = Path(__file__).parent / "diagrams"
diagram_files = [
    ("01_overview_30k.png", "Architecture: 30,000ft Overview"),
    ("02_pipeline_2k.png", "Architecture: Pipeline Detail"),
    ("03_backend_2k.png", "Architecture: Backend Detail"),
    ("04_frontend_2k.png", "Architecture: Frontend Detail"),
    ("05_entity_model_2k.png", "Architecture: Entity Model Detail"),
]
for fname, title in diagram_files:
    fpath = diagram_dir / fname
    if fpath.exists():
        s = add_slide()
        add_text(s, 0.5, 0.15, 12, 0.35, title, 20, ACCENT, bold=True)
        pic = s.shapes.add_picture(str(fpath), Inches(0.3), Inches(0.55), width=Inches(12.7))

# Save
out = "undata_presentation.pptx"
prs.save(out)
print(f"Created: {out}")
